"""RAG utilities for document discovery, indexing, and chat.

Handles optional dependencies gracefully, adds concise docstrings, and
normalizes typing to play well with Pylance.
"""

from __future__ import annotations

import contextlib
import csv
import hashlib
import json
import os
import sys
import time
from collections.abc import Callable
from dataclasses import dataclass
from html import unescape
from html.parser import HTMLParser
from typing import TYPE_CHECKING, Any

import numpy as np
from openai import OpenAI

# Optional deps (guarded imports)
try:  # python-docx
    import docx  # type: ignore
except Exception:  # pragma: no cover - optional
    docx = None  # type: ignore

try:  # openpyxl
    import openpyxl  # type: ignore
except Exception:  # pragma: no cover - optional
    openpyxl = None  # type: ignore

try:  # python-pptx
    import pptx  # type: ignore
except Exception:  # pragma: no cover - optional
    pptx = None  # type: ignore

try:  # pypdf
    from pypdf import PdfReader  # type: ignore
except Exception:  # pragma: no cover - optional
    PdfReader = None  # type: ignore

try:  # pdfminer
    from pdfminer.high_level import extract_text as pdfminer_extract_text  # type: ignore

    HAS_PDFMINER = True
except Exception:  # pragma: no cover - optional
    pdfminer_extract_text = None  # type: ignore
    HAS_PDFMINER = False

try:  # striprtf
    from striprtf.striprtf import rtf_to_text  # type: ignore
except Exception:  # pragma: no cover - optional

    def rtf_to_text(_: str) -> str:  # type: ignore
        return ""


APP_NAME = "DocChatbot"
DEFAULT_EMBED_MODEL = "text-embedding-3-small"
DEFAULT_CHAT_MODEL = "gpt-4o-mini"


def base_dir() -> str:
    """Return the application base directory.

    Uses the executable path when bundled (``sys.frozen``), else this file's folder.
    """
    if getattr(sys, 'frozen', False):
        return os.path.dirname(sys.executable)
    return os.path.dirname(os.path.abspath(__file__))


def paths() -> dict[str, str]:
    """Return important application paths as a mapping."""
    root = base_dir()
    return {
        "root": root,
        "docs": os.path.join(root, "docs"),
        "storage": os.path.join(root, "storage"),
        "index": os.path.join(root, "storage", "index.json"),
        "config": os.path.join(root, "storage", "config.json"),
        "api_key": os.path.join(root, "api_key.txt"),
        "log": os.path.join(root, "storage", "app.log"),
    }


def ensure_dirs() -> None:
    """Ensure required directories exist."""
    p = paths()
    os.makedirs(p["docs"], exist_ok=True)
    os.makedirs(p["storage"], exist_ok=True)


def read_api_key() -> str | None:
    """Read the API key from ``api_key.txt`` (first line) if present."""
    p = paths()["api_key"]
    if not os.path.exists(p):
        return None
    with open(p, encoding="utf-8") as f:
        key = f.readline().strip()
        return key or None


def sha1_file(path: str) -> str:
    """Compute a file's SHA-1 hash in hex."""
    h = hashlib.sha1()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()


def chunk_text(text: str, chunk_chars: int = 1200, overlap: int = 200) -> list[str]:
    """Split text into overlapping character chunks for embedding."""
    text = text.replace("\r\n", "\n")
    parts = [p.strip() for p in text.split("\n\n") if p.strip()]
    chunks: list[str] = []
    buf = ""
    for part in parts:
        if len(buf) + len(part) + 2 <= chunk_chars:
            buf = (buf + "\n\n" + part) if buf else part
        else:
            if buf:
                chunks.append(buf)
            # if a single part is huge, slice it
            if len(part) > chunk_chars:
                i = 0
                while i < len(part):
                    chunks.append(part[i : i + chunk_chars])
                    i += chunk_chars - overlap
                buf = ""
            else:
                buf = part
    if buf:
        chunks.append(buf)

    # Add overlaps across boundaries
    with_overlap: list[str] = []
    for i, c in enumerate(chunks):
        if i == 0 or overlap <= 0:
            with_overlap.append(c)
        else:
            prefix = chunks[i - 1][-overlap:]
            with_overlap.append(prefix + c)
    return with_overlap


def extract_pdf_text(pdf_path: str) -> str:
    """Extract text from a PDF using pypdf, with pdfminer fallback."""
    try:
        text = ""
        if PdfReader is not None:
            reader = PdfReader(pdf_path)  # type: ignore[misc]
            texts: list[str] = []
            for page in reader.pages:  # type: ignore[union-attr]
                t = page.extract_text() or ""
                texts.append(t)
            text = "\n\n".join(texts)
    except Exception:
        text = ""

    if (not text or not text.strip()) and HAS_PDFMINER and pdfminer_extract_text is not None:
        with contextlib.suppress(Exception):
            text = pdfminer_extract_text(pdf_path) or ""
    return text or ""


def extract_docx_text(path: str) -> str:
    """Extract text from a .docx file using python-docx (optional)."""
    try:
        if docx is None:
            return ""
        d = docx.Document(path)  # type: ignore[operator]
        return "\n\n".join(p.text for p in d.paragraphs if p.text)
    except Exception:
        return ""


def extract_plain_text(path: str) -> str:
    try:
        with open(path, encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception:
        return ""


class _HTMLTextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self._ignore_stack: list[str] = []
        self._parts: list[str] = []

    def handle_starttag(self, tag: str, attrs):  # noqa: ANN001 - HTMLParser signature
        if tag in {"script", "style"}:
            self._ignore_stack.append(tag)
        elif tag in {
            "p",
            "br",
            "div",
            "li",
            "section",
            "article",
            "h1",
            "h2",
            "h3",
            "h4",
            "h5",
            "h6",
        }:
            self._parts.append("\n")

    def handle_endtag(self, tag: str):  # noqa: ANN001 - HTMLParser signature
        if self._ignore_stack and self._ignore_stack[-1] == tag:
            self._ignore_stack.pop()

    def handle_data(self, data: str):  # noqa: ANN001 - HTMLParser signature
        if not self._ignore_stack:
            text = unescape(data).strip()
            if text:
                self._parts.append(text)

    def get_text(self) -> str:
        joined = " ".join(self._parts)
        lines = [line.strip() for line in joined.replace("\r\n", "\n").split("\n")]
        return "\n".join([ln for ln in lines if ln])


def extract_html_text(path: str) -> str:
    """Extract readable text from an HTML file."""
    try:
        with open(path, encoding="utf-8", errors="ignore") as f:
            html = f.read()
        parser = _HTMLTextExtractor()
        parser.feed(html)
        return parser.get_text()
    except Exception:
        return ""


def extract_csv_text(path: str, limit_rows: int = 200) -> str:
    """Extract a quick text preview from a CSV (limited rows)."""
    try:
        lines: list[str] = []
        with open(path, encoding="utf-8", errors="ignore", newline="") as f:
            reader = csv.reader(f)
            for i, row in enumerate(reader):
                lines.append(", ".join([str(x) for x in row]))
                if i + 1 >= limit_rows:
                    break
        return "\n".join(lines)
    except Exception:
        return extract_plain_text(path)


def extract_pptx_text(path: str) -> str:
    """Extract text from a PowerPoint (.pptx) using python-pptx (optional)."""
    try:
        if pptx is None:
            return ""
        prs = pptx.Presentation(path)  # type: ignore[operator]
        texts: list[str] = []
        for slide in prs.slides:  # type: ignore[union-attr]
            for shape in slide.shapes:  # type: ignore[union-attr]
                t = getattr(shape, "text", None)
                if t:
                    texts.append(str(t))
        return "\n\n".join(texts)
    except Exception:
        return ""


def extract_xlsx_text(path: str, limit_cells: int = 5000) -> str:
    """Extract a rough text preview from an Excel workbook (optional)."""
    try:
        if openpyxl is None:
            return ""
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)  # type: ignore[operator]
        out: list[str] = []
        count = 0
        for ws in wb.worksheets:  # type: ignore[union-attr]
            out.append(f"# Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):  # type: ignore[union-attr]
                if row is None:
                    continue
                out.append(", ".join(["" if c is None else str(c) for c in row]))
                count += len(row)
                if count >= limit_cells:
                    break
            if count >= limit_cells:
                break
        return "\n".join(out)
    except Exception:
        return ""


def extract_rtf_text(path: str) -> str:
    """Extract text from an RTF file (optional via striprtf)."""
    try:
        with open(path, encoding="utf-8", errors="ignore") as f:
            data = f.read()
        return rtf_to_text(data)
    except Exception:
        return ""


SUPPORTED_EXTS = {".pdf", ".docx", ".txt", ".md", ".csv", ".html", ".htm", ".pptx", ".xlsx", ".rtf"}


def list_docs() -> list[str]:
    """Return a sorted list of document file paths under ``docs/``."""
    d = paths()["docs"]
    out: list[str] = []
    for root, _, files in os.walk(d):
        for fn in files:
            ext = os.path.splitext(fn)[1].lower()
            if ext in SUPPORTED_EXTS:
                out.append(os.path.join(root, fn))
    return sorted(out)


@dataclass
class IndexItem:
    """Data for an embedded chunk of a document."""

    id: str
    doc_path: str
    doc_name: str
    mtime: float
    chunk_index: int
    text: str
    embedding: list[float]


def _load_index_raw(index_path: str) -> dict:
    """Load the raw JSON index mapping from disk or return empty mapping."""
    if not os.path.exists(index_path):
        return {}
    try:
        with open(index_path, encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}


def _save_index_raw(index_path: str, data: dict) -> None:
    """Atomically save the JSON index mapping to disk."""
    tmp = index_path + ".tmp"
    with open(tmp, "w", encoding="utf-8") as f:
        json.dump(data, f)
    os.replace(tmp, index_path)


def _norm_rows(mat: np.ndarray) -> np.ndarray:
    """Row-normalize a matrix to unit length (L2)."""
    norms = np.linalg.norm(mat, axis=1, keepdims=True) + 1e-12
    return mat / norms


def create_client(api_key: str) -> OpenAI:
    """Create an OpenAI client bound to the provided API key."""
    return OpenAI(api_key=api_key)


def embed_texts(client: OpenAI, texts: list[str], embed_model: str) -> list[list[float]]:
    """Embed a list of texts using the specified embedding model."""
    out: list[list[float]] = []
    B = 64
    for i in range(0, len(texts), B):
        batch = texts[i : i + B]
        resp = client.embeddings.create(model=embed_model, input=batch)
        for item in resp.data:
            out.append(item.embedding)
    return out


ProgressCb = Callable[[str, int, int], None] | None


def _extract_text_for_path(path_any: str) -> str:
    """Dispatch to the appropriate text extractor based on file extension."""
    ext = os.path.splitext(path_any)[1].lower()
    if ext == ".pdf":
        return extract_pdf_text(path_any)
    if ext == ".docx":
        return extract_docx_text(path_any)
    if ext in {".txt", ".md"}:
        return extract_plain_text(path_any)
    if ext in {".html", ".htm"}:
        return extract_html_text(path_any)
    if ext == ".csv":
        return extract_csv_text(path_any)
    if ext == ".pptx":
        return extract_pptx_text(path_any)
    if ext == ".xlsx":
        return extract_xlsx_text(path_any)
    if ext == ".rtf":
        return extract_rtf_text(path_any)
    return ""


def build_or_update_index(
    client: OpenAI,
    status_cb=None,
    progress_cb: ProgressCb = None,
    *,
    embed_model: str = DEFAULT_EMBED_MODEL,
    chunk_chars: int = 1200,
    overlap: int = 200,
    force_full: bool = False,
) -> tuple[list[IndexItem], np.ndarray]:
    """Build or incrementally update the embeddings index.

    Returns a tuple of (items, matrix) where matrix rows are unit-normalized.
    """
    p = paths()
    ensure_dirs()
    index_path = p["index"]
    raw = _load_index_raw(index_path)
    items: list[IndexItem] = []

    # If embedding model changed, force full rebuild
    if raw and raw.get("embed_model") and raw.get("embed_model") != embed_model:
        force_full = True

    # map of existing entries by (doc_path, chunk_index)
    existing: dict[tuple[str, int], dict] = {}
    if raw.get("items"):
        for it in raw["items"]:
            key = (it["doc_path"], it["chunk_index"])
            existing[key] = it

    docs = list_docs()
    to_add: list[IndexItem] = []
    total_new_chunks = 0

    # Pre-scan to compute chunks count for progress
    prescan: list[tuple[str, float, str, list[str]]] = []  # (path, mtime, doc_name, chunks)
    for path_any in docs:
        mtime = os.path.getmtime(path_any)
        doc_name = os.path.basename(path_any)
        # Decide whether to (re)ingest this file
        needs_reingest = bool(force_full)
        if not force_full:
            needs_reingest = True
            for key, it in list(existing.items()):
                if key[0] == path_any and it.get("mtime", 0) >= mtime:
                    needs_reingest = False
                    break
        if needs_reingest:
            if status_cb:
                status_cb(f"Reading {doc_name}…")
            text = _extract_text_for_path(path_any)
            if not text.strip():
                continue
            chunks = chunk_text(text, chunk_chars=chunk_chars, overlap=overlap)
            prescan.append((path_any, mtime, doc_name, chunks))
            total_new_chunks += len(chunks)
        else:
            prescan.append((path_any, mtime, doc_name, []))

    # Embed in batches across all docs, updating progress
    done_chunks = 0
    for path_any, mtime, doc_name, chunks in prescan:
        if chunks:
            if status_cb:
                status_cb(f"Embedding {doc_name} ({len(chunks)} chunks)…")
            if progress_cb and total_new_chunks > 0:
                progress_cb("embedding", done_chunks, total_new_chunks)
            embs = embed_texts(client, chunks, embed_model)
            for ci, (t, e) in enumerate(zip(chunks, embs, strict=False)):
                uid = hashlib.sha1(f"{path_any}:{mtime}:{ci}".encode()).hexdigest()
                to_add.append(
                    IndexItem(
                        id=uid,
                        doc_path=path_any,
                        doc_name=doc_name,
                        mtime=mtime,
                        chunk_index=ci,
                        text=t,
                        embedding=e,
                    )
                )
            done_chunks += len(chunks)
            if progress_cb and total_new_chunks > 0:
                progress_cb("embedding", done_chunks, total_new_chunks)

    # Rebuild final list: include up-to-date existing + new
    kept: list[IndexItem] = []
    for it in raw.get("items", []):
        # Keep only if file still exists and not superseded
        if (
            not force_full
            and os.path.exists(it["doc_path"])
            and (it["doc_path"], it["chunk_index"])
            not in [(x.doc_path, x.chunk_index) for x in to_add]
        ):
            kept.append(
                IndexItem(
                    id=it["id"],
                    doc_path=it["doc_path"],
                    doc_name=it["doc_name"],
                    mtime=it.get("mtime", 0.0),
                    chunk_index=it["chunk_index"],
                    text=it["text"],
                    embedding=it["embedding"],
                )
            )

    items = kept + to_add

    # Save
    raw_out = {
        "version": 1,
        "built_at": time.time(),
        "embed_model": embed_model,
        "embed_dim": (len(items[0].embedding) if items else 0),
        "items": [
            {
                "id": it.id,
                "doc_path": it.doc_path,
                "doc_name": it.doc_name,
                "mtime": it.mtime,
                "chunk_index": it.chunk_index,
                "text": it.text,
                "embedding": it.embedding,
            }
            for it in items
        ],
    }
    _save_index_raw(index_path, raw_out)

    # Build matrix
    if items:
        mat = np.array([it.embedding for it in items], dtype=np.float32)
        mat = _norm_rows(mat)
    else:
        dim = raw.get("embed_dim", 0) if isinstance(raw, dict) else 0
        mat = np.zeros((0, dim), dtype=np.float32)

    if status_cb:
        status_cb(f"Index ready ({len(items)} chunks)")

    return items, mat


def load_index_cached() -> tuple[list[IndexItem], np.ndarray]:
    p = paths()
    raw = _load_index_raw(p["index"]) or {}
    items_raw = raw.get("items", [])
    items: list[IndexItem] = [
        IndexItem(
            id=it["id"],
            doc_path=it["doc_path"],
            doc_name=it["doc_name"],
            mtime=it.get("mtime", 0.0),
            chunk_index=it["chunk_index"],
            text=it["text"],
            embedding=it["embedding"],
        )
        for it in items_raw
    ]
    if items:
        mat = np.array([it.embedding for it in items], dtype=np.float32)
        mat = _norm_rows(mat)
    else:
        dim = raw.get("embed_dim", 0)
        mat = np.zeros((0, dim), dtype=np.float32)
    return items, mat


def retrieve(
    client: OpenAI,
    items: list[IndexItem],
    mat: np.ndarray,
    query: str,
    *,
    top_k: int = 4,
    embed_model: str = DEFAULT_EMBED_MODEL,
    mmr: bool = True,
    mmr_lambda: float = 0.5,
) -> list[tuple[IndexItem, float]]:
    if not items:
        return []
    q_emb = client.embeddings.create(model=embed_model, input=[query]).data[0].embedding
    q = np.array(q_emb, dtype=np.float32)
    q = q / (np.linalg.norm(q) + 1e-12)
    sims = mat @ q  # cosine similarity

    if not mmr:
        idx = np.argsort(-sims)[:top_k]
        return [(items[int(i)], float(sims[int(i)])) for i in idx]

    # MMR selection
    selected: list[int] = []
    candidates = list(np.argsort(-sims))  # sorted best-first
    if not candidates:
        return []
    selected.append(int(candidates.pop(0)))
    # Precompute norms and similarities between all docs for diversity
    # Using cosine sim via dot products; mat rows are unit norm.
    for _ in range(top_k - 1):
        best_score = -1e9
        best_idx = None
        for c in candidates:
            # relevance term
            rel = sims[int(c)]
            # diversity term: max similarity to already selected
            div = 0.0
            for s in selected:
                div = max(div, float(mat[int(c)] @ mat[int(s)]))
            score = mmr_lambda * rel - (1 - mmr_lambda) * div
            if score > best_score:
                best_score = score
                best_idx = int(c)
        if best_idx is None:
            break
        selected.append(best_idx)
        candidates.remove(best_idx)
        if not candidates:
            break
    return [(items[i], float(sims[i])) for i in selected]


if TYPE_CHECKING:
    from openai.types.chat import ChatCompletionMessageParam as _Msg
else:  # At runtime, keep it simple to avoid heavy imports
    _Msg = dict[str, str]  # type: ignore[misc, assignment]


def build_prompt(query: str, retrieved: list[IndexItem]) -> list[_Msg]:
    context_blocks = []
    for i, it in enumerate(retrieved, start=1):
        tag = f"[{i}] {it.doc_name} (chunk {it.chunk_index + 1})"
        context_blocks.append(f"{tag}\n{it.text}")

    context = "\n\n".join(context_blocks) if context_blocks else "(no relevant context found)"

    system = (
        "You are a helpful assistant that answers questions using only the provided context. "
        "If the answer is not present in the context, say you don't know and suggest adding more documents. "  # noqa: E501
        "Cite sources inline like [1], [2] and list sources at the end. Be concise."
    )

    user = (
        f"User question:\n{query}\n\n"
        f"Context:\n{context}\n\n"
        f"Instructions: Base your answer strictly on the context above."
    )

    return [
        {"role": "system", "content": system},  # type: ignore[typeddict-item]
        {"role": "user", "content": user},  # type: ignore[typeddict-item]
    ]


def chat_answer(
    client: OpenAI,
    query: str,
    items: list[IndexItem],
    mat: np.ndarray,
    *,
    chat_model: str = DEFAULT_CHAT_MODEL,
    embed_model: str = DEFAULT_EMBED_MODEL,
    top_k: int = 4,
    mmr: bool = True,
    mmr_lambda: float = 0.5,
    min_relevance: float = 0.12,
) -> tuple[str, list[str]]:
    retrieved_scored = retrieve(
        client,
        items,
        mat,
        query,
        top_k=top_k,
        embed_model=embed_model,
        mmr=mmr,
        mmr_lambda=mmr_lambda,
    )
    if not retrieved_scored:
        raise RuntimeError("No relevant information found in your documents for this question.")
    best_score = max(s for _, s in retrieved_scored)
    if best_score < min_relevance:
        raise RuntimeError("No relevant information found in your documents for this question.")

    retrieved_items = [it for it, _ in retrieved_scored]
    msgs = build_prompt(query, retrieved_items)
    resp = client.chat.completions.create(
        model=chat_model,
        messages=msgs,  # type: ignore[arg-type]
        temperature=0.2,
    )
    answer = resp.choices[0].message.content or ""
    srcs = [f"[{i + 1}] {it.doc_name}" for i, it in enumerate(retrieved_items)]
    return answer, srcs
